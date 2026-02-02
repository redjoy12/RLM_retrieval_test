"""Office document parsers (DOCX, XLSX, PPTX)."""

import asyncio
from pathlib import Path
from typing import List

import aiofiles
from docx import Document as DocxDocument
from openpyxl import load_workbook
from pptx import Presentation

from rlm.documents.models import (
    ContentChunk,
    DocumentContent,
    DocumentFormat,
    DocumentMetadata,
    DocumentStructure,
    ParseOptions,
)
from rlm.documents.parsers.base import DocumentParser
from rlm.documents.parsers.text_parser import ParseError


class DocxParser(DocumentParser):
    """Parser for Word documents (.docx)."""

    SUPPORTED_FORMATS = {DocumentFormat.DOCX}

    @property
    def name(self) -> str:
        return "docx"

    @property
    def priority(self) -> int:
        return 100

    def supports(self, format_type: DocumentFormat) -> bool:
        return format_type in self.SUPPORTED_FORMATS

    async def parse(
        self, file_path: Path, options: ParseOptions
    ) -> DocumentContent:
        loop = asyncio.get_event_loop()
        return await loop.run_in_executor(
            None, self._parse_sync, file_path, options
        )

    def _parse_sync(self, file_path: Path, options: ParseOptions) -> DocumentContent:
        try:
            doc = DocxDocument(file_path)

            # Extract metadata
            metadata = self._extract_metadata(doc)

            # Extract text from paragraphs
            paragraphs = []
            for para in doc.paragraphs:
                if para.text.strip():
                    paragraphs.append(para.text)

            # Extract text from tables if requested
            if options.extract_tables:
                for table in doc.tables:
                    table_text = self._table_to_text(table)
                    if table_text:
                        paragraphs.append(f"[TABLE]\n{table_text}")

            full_text = "\n\n".join(paragraphs)

            # Extract structure
            structure = DocumentStructure(
                title=metadata.title,
                headings=self._extract_headings(doc),
            )

            # Create chunks
            chunks = self._chunk_content(full_text, options)

            return DocumentContent(
                raw_text=full_text,
                cleaned_text=full_text,
                chunks=chunks,
                structure=structure,
            )

        except Exception as e:
            raise ParseError(f"Failed to parse DOCX: {e}") from e

    def _extract_metadata(self, doc: DocxDocument) -> DocumentMetadata:
        """Extract metadata from Word document."""
        core_props = doc.core_properties

        return DocumentMetadata(
            title=core_props.title,
            author=core_props.author,
            subject=core_props.subject,
            created_date=core_props.created,
            modified_date=core_props.modified,
            word_count=len(doc.paragraphs),
        )

    def _extract_headings(self, doc: DocxDocument) -> List[str]:
        """Extract headings from document."""
        headings = []
        for para in doc.paragraphs:
            if para.style.name.startswith("Heading"):
                headings.append(para.text)
        return headings

    def _table_to_text(self, table) -> str:
        """Convert table to text."""
        rows = []
        for row in table.rows:
            row_text = " | ".join(cell.text.strip() for cell in row.cells)
            rows.append(row_text)
        return "\n".join(rows)

    def _chunk_content(self, text: str, options: ParseOptions) -> List[ContentChunk]:
        """Split content into chunks."""
        chunks = []
        paragraphs = text.split("\n\n")
        current_chunk = ""
        char_pos = 0
        chunk_idx = 0
        chunk_size = options.chunk_size
        overlap = options.chunk_overlap

        for para in paragraphs:
            para = para.strip()
            if not para:
                continue

            para_with_newline = para + "\n\n"

            if len(current_chunk) + len(para_with_newline) > chunk_size:
                if current_chunk:
                    chunks.append(
                        ContentChunk(
                            index=chunk_idx,
                            content=current_chunk.strip(),
                            start_char=char_pos,
                            end_char=char_pos + len(current_chunk),
                        )
                    )
                    chunk_idx += 1
                    char_pos += len(current_chunk) - overlap

                if overlap > 0 and current_chunk:
                    current_chunk = current_chunk[-overlap:] + para_with_newline
                else:
                    current_chunk = para_with_newline
            else:
                current_chunk += para_with_newline

        if current_chunk.strip():
            chunks.append(
                ContentChunk(
                    index=chunk_idx,
                    content=current_chunk.strip(),
                    start_char=char_pos,
                    end_char=char_pos + len(current_chunk),
                )
            )

        return chunks


class XlsxParser(DocumentParser):
    """Parser for Excel spreadsheets (.xlsx)."""

    SUPPORTED_FORMATS = {DocumentFormat.XLSX}

    @property
    def name(self) -> str:
        return "xlsx"

    @property
    def priority(self) -> int:
        return 100

    def supports(self, format_type: DocumentFormat) -> bool:
        return format_type in self.SUPPORTED_FORMATS

    async def parse(
        self, file_path: Path, options: ParseOptions
    ) -> DocumentContent:
        loop = asyncio.get_event_loop()
        return await loop.run_in_executor(
            None, self._parse_sync, file_path, options
        )

    def _parse_sync(self, file_path: Path, options: ParseOptions) -> DocumentContent:
        try:
            wb = load_workbook(file_path, data_only=True)

            sheets_text = []
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                sheet_text = [f"[SHEET: {sheet_name}]"]

                for row in sheet.iter_rows():
                    row_values = []
                    for cell in row:
                        value = cell.value
                        if value is not None:
                            row_values.append(str(value))

                    if row_values:
                        sheet_text.append(" | ".join(row_values))

                if len(sheet_text) > 1:  # More than just the header
                    sheets_text.append("\n".join(sheet_text))

            full_text = "\n\n".join(sheets_text)

            # Single chunk for spreadsheet
            chunks = [
                ContentChunk(
                    index=0,
                    content=full_text,
                    start_char=0,
                    end_char=len(full_text),
                )
            ]

            return DocumentContent(
                raw_text=full_text,
                cleaned_text=full_text,
                chunks=chunks,
                structure=DocumentStructure(),
            )

        except Exception as e:
            raise ParseError(f"Failed to parse XLSX: {e}") from e


class PptxParser(DocumentParser):
    """Parser for PowerPoint presentations (.pptx)."""

    SUPPORTED_FORMATS = {DocumentFormat.PPTX}

    @property
    def name(self) -> str:
        return "pptx"

    @property
    def priority(self) -> int:
        return 100

    def supports(self, format_type: DocumentFormat) -> bool:
        return format_type in self.SUPPORTED_FORMATS

    async def parse(
        self, file_path: Path, options: ParseOptions
    ) -> DocumentContent:
        loop = asyncio.get_event_loop()
        return await loop.run_in_executor(
            None, self._parse_sync, file_path, options
        )

    def _parse_sync(self, file_path: Path, options: ParseOptions) -> DocumentContent:
        try:
            prs = Presentation(file_path)

            slides_text = []
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_content = [f"[SLIDE {slide_num}]"]

                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_content.append(shape.text.strip())

                if len(slide_content) > 1:
                    slides_text.append("\n".join(slide_content))

            full_text = "\n\n".join(slides_text)

            # Create chunks by slide
            chunks = []
            char_pos = 0
            for idx, slide_text in enumerate(slides_text):
                chunks.append(
                    ContentChunk(
                        index=idx,
                        content=slide_text,
                        start_char=char_pos,
                        end_char=char_pos + len(slide_text),
                    )
                )
                char_pos += len(slide_text)

            return DocumentContent(
                raw_text=full_text,
                cleaned_text=full_text,
                chunks=chunks,
                structure=DocumentStructure(),
            )

        except Exception as e:
            raise ParseError(f"Failed to parse PPTX: {e}") from e
